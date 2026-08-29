import os
# Fixes the PyTorch HuggingFace silent crash
os.environ["KMP_DUPLICATE_LIB_OK"] = "TRUE"
# Fixes the Protobuf websocket silent crash
os.environ["PROTOCOL_BUFFERS_PYTHON_IMPLEMENTATION"] = "python" 

import tempfile
import warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)

import streamlit as st
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation

from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain_core.prompts import PromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_core.output_parsers import StrOutputParser

# ==========================================
# HARDCODED API KEY (Replace with your actual key)
# ==========================================
GOOGLE_API_KEY = st.secrets["GOOGLE_API_KEY"]

# 1. Setup Streamlit Page
st.set_page_config(page_title="Universal Multi-File RAG", page_icon="🗂️", layout="wide")
st.title("🗂️ Universal Multi-File RAG Assistant")
st.caption("Supports PDF, Word (.docx), PowerPoint (.pptx), Excel (.xlsx), CSV, and Text files.")

# 2. Universal File Loader Function (Windows Safe)
def extract_text_from_file(file_path: str, filename: str) -> list[Document]:
    """Extracts raw text and converts it into LangChain Document format."""
    ext = os.path.splitext(filename)[1].lower()
    docs = []

    try:
        if ext == ".pdf":
            loader = PyPDFLoader(file_path)
            docs = loader.load()
            for doc in docs:
                doc.metadata["source"] = filename

        elif ext in [".docx", ".doc"]:
            # Context manager drops the file lock immediately
            with open(file_path, "rb") as f:
                docx_obj = DocxDocument(f)
                full_text = "\n".join([p.text for p in docx_obj.paragraphs if p.text.strip()])
            docs = [Document(page_content=full_text, metadata={"source": filename})]

        elif ext in [".pptx", ".ppt"]:
            with open(file_path, "rb") as f:
                prs = Presentation(f)
                slide_texts = []
                for i, slide in enumerate(prs.slides):
                    slide_content = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                if paragraph.text.strip():
                                    slide_content.append(paragraph.text.strip())
                    if slide_content:
                        slide_texts.append(f"[Slide {i+1}]\n" + "\n".join(slide_content))
                full_text = "\n\n".join(slide_texts)
            docs = [Document(page_content=full_text, metadata={"source": filename})]

        elif ext in [".xlsx", ".xls"]:
            # sheet_name=None reads all sheets and automatically closes the file
            sheets = pd.read_excel(file_path, sheet_name=None)
            sheets_text = []
            for sheet_name, df in sheets.items():
                sheets_text.append(f"[Sheet: {sheet_name}]\n" + df.to_string(index=False))
            full_text = "\n\n".join(sheets_text)
            docs = [Document(page_content=full_text, metadata={"source": filename})]

        elif ext == ".csv":
            df = pd.read_csv(file_path)
            docs = [Document(page_content=df.to_string(index=False), metadata={"source": filename})]

        elif ext in [".txt", ".md", ".json", ".log"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            docs = [Document(page_content=text, metadata={"source": filename})]

    except Exception as e:
        st.warning(f"Could not read {filename}: {e}")

    return docs

# 3. Cache Database Creation for Multiple Files
@st.cache_resource(show_spinner=False)
def build_vectorstore(uploaded_files):
    all_raw_docs = []
    
    for uploaded_file in uploaded_files:
        ext = os.path.splitext(uploaded_file.name)[1]
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            tmp_file_path = tmp_file.name

        docs = extract_text_from_file(tmp_file_path, uploaded_file.name)
        all_raw_docs.extend(docs)
        
        # Windows-Safe Temp File Deletion
        try:
            if os.path.exists(tmp_file_path):
                os.remove(tmp_file_path)
        except Exception:
            pass # Ignore if a background task still momentarily holds the file

    # Chunking
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    chunks = text_splitter.split_documents(all_raw_docs)

    # Embed & Store
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    vectorstore = Chroma.from_documents(documents=chunks, embedding=embeddings)
    
    return vectorstore

def format_docs(docs):
    return "\n\n".join(f"[Source: {doc.metadata.get('source', 'Unknown')}]\n{doc.page_content}" for doc in docs)

# 4. Multi-File Uploader in Sidebar
st.sidebar.header("📁 Knowledge Base")
uploaded_files = st.sidebar.file_uploader(
    "Upload multiple files",
    type=["pdf", "docx", "pptx", "xlsx", "xls", "csv", "txt", "md"],
    accept_multiple_files=True
)

if uploaded_files:
    with st.spinner(f"Ingesting & indexing {len(uploaded_files)} file(s)..."):
        vectorstore = build_vectorstore(uploaded_files)
        retriever = vectorstore.as_retriever(search_kwargs={"k": 20})
        st.sidebar.success(f"Indexed {len(uploaded_files)} file(s) successfully!")

    # ==========================================
    # 5. Database Viewer Section
    # ==========================================
    st.sidebar.markdown("---")
    st.sidebar.header("🗄️ Database Inspector")
    
    if st.sidebar.checkbox("Show Database Records"):
        st.subheader("📊 Vector Database Contents")
        
        # Retrieve all stored vectors/chunks directly from Chroma
        raw_db_data = vectorstore.get()
        total_chunks = len(raw_db_data.get("ids", []))
        
        col1, col2 = st.columns(2)
        col1.metric("Total Files Uploaded", len(uploaded_files))
        col2.metric("Total Chunks in Vector DB", total_chunks)
        
        if total_chunks > 0:
            records = []
            for i, chunk_id in enumerate(raw_db_data["ids"]):
                meta = raw_db_data["metadatas"][i] if raw_db_data["metadatas"] else {}
                content = raw_db_data["documents"][i]
                records.append({
                    "Index": i + 1,
                    "Chunk ID": chunk_id,
                    "Source": meta.get("source", "N/A"),
                    "Preview": content[:120].replace("\n", " ") + "...",
                    "Full Content": content
                })
            
            df = pd.DataFrame(records)
            st.dataframe(df[["Index", "Source", "Preview", "Chunk ID"]], use_container_width=True)
            
            with st.expander("🔍 Inspect Specific Chunk"):
                selected_idx = st.selectbox("Select chunk to view complete text:", df["Index"])
                selected_row = df[df["Index"] == selected_idx].iloc[0]
                st.write(f"**Source Document:** `{selected_row['Source']}`")
                st.write(f"**Chunk ID:** `{selected_row['Chunk ID']}`")
                st.text_area("Full Extracted Text", selected_row["Full Content"], height=200)
        
        st.markdown("---")

    # ==========================================
    # 6. RAG Pipeline Setup
    # ==========================================
    llm = ChatGoogleGenerativeAI(
        model="gemini-2.5-flash", 
        temperature=0, 
        google_api_key=GOOGLE_API_KEY
    )
    
    prompt = PromptTemplate.from_template(
        """Answer the question using ONLY the provided context. 
        Cite the source filename for each fact whenever possible.
        If you don't know the answer based on the context, say so. Do not make anything up.
        
        Context:
        {context}
        
        Question: {question}
        
        Answer:"""
    )
    
    rag_chain = (
        {"context": retriever | format_docs, "question": RunnablePassthrough()}
        | prompt
        | llm
        | StrOutputParser()
    )

    # ==========================================
    # 7. Chat Interface
    # ==========================================
    if "messages" not in st.session_state:
        st.session_state.messages = []

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])

    if question := st.chat_input("Ask anything across all uploaded documents..."):
        st.session_state.messages.append({"role": "user", "content": question})
        with st.chat_message("user"):
            st.markdown(question)

        with st.chat_message("assistant"):
            with st.spinner("Searching across all files & generating answer..."):
                response = rag_chain.invoke(question)
                st.markdown(response)
        
        st.session_state.messages.append({"role": "assistant", "content": response})

else:
    st.info("👈 Upload one or more documents in the sidebar to start.")